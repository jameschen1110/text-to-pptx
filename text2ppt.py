from ollama import ChatResponse, chat
from pptx import Presentation
import json

# import needed functions

FORMAT = [
    {"title": "<page 1 title>", "content": "<page 1 long content.>"},
    {"title": "<page 2 title>", "content": "<page 2 long content.\nline2>"},
]


def apply(msg1, msg2):  # asking function
    response: ChatResponse = chat(
        model="deepseek-r1:7b",
        messages=[
            {
                "role": "system",
                # give order as system(weigh stronger than user)
                "content": f"""\
output strictly and directly as following JSON format without any other text, for example:
```json
{json.dumps(FORMAT)}
```
for breakline, using "\\n", do not break a new line directly
""",  # setting format
            },
            {
                "role": "user",
                # give order as user
                "content": "generate a presentation about '" + msg1 + "' within " + msg2 + " pages",
                # asking in format, msg1:theme, msg2:pages
            },
        ],
        stream=True,
    )
    return response


# main start
print('Type "/exit" to exit.')
while True:
    text1 = input("Presentation theme: ").strip()
    if text1 == "/exit":
        break
    text2 = input("Pages required: ").strip()
    if text2 == "/exit":
        break
    # input things required
    total_ans = ""  # initialize
    stream = apply(text1, text2)  # asking
    for chunk in stream:  # receive answer per token(word)
        mes = chunk["message"]["content"]
        print(mes, end="")  # output in command line
        total_ans = total_ans + mes  # add answer into total_ans
    # total_ans.replace("\\n", "\n")  # replace \n into real enter(switch line)

    print("\n===== Generate END =====")

    data = json.loads(total_ans.split("```json")[1].strip().removesuffix("```"))  # convert into json format
    prs = Presentation()  # create a presentation
    title_slide_layout = prs.slide_layouts[1]  # set presentation format
    for page in data:
        slide = prs.slides.add_slide(title_slide_layout)  # add a slide
        title = slide.shapes.title  # title setting
        content = slide.placeholders[1]  # content setting

        title.text = page["title"]
        content.text = page["content"]

    print("===== ALL END =====")
    prs.save(text1 + ".pptx")
    # save file in theme name
